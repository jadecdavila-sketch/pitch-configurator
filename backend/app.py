from flask import Flask, request, send_file, jsonify
from flask_cors import CORS
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from copy import deepcopy
import io
import os
from datetime import datetime
from lxml import etree
from dotenv import load_dotenv
import google.generativeai as genai

# Load environment variables
load_dotenv()

app = Flask(__name__)
CORS(app)

# Configure Gemini API
GEMINI_API_KEY = os.getenv('GEMINI_API_KEY')
if GEMINI_API_KEY:
    genai.configure(api_key=GEMINI_API_KEY)
    print("Gemini API configured successfully")
else:
    print("WARNING: GEMINI_API_KEY not found in environment variables")

# Recipe to slide mapping (matches frontend recipeSlideMapping.ts)
RECIPE_TO_SLIDE_MAP = {
    # Individual Contributor Courses (slides 3-7)
    'day1-onboarding': 3,
    '30-60-90-onboarding': 4,
    'critical-thinking': 5,
    'navigating-matrix': 6,
    'decoding-business': 7,

    # Manager Courses (slides 9-16)
    'ascend-leadership': 9,
    'guiding-performance': 10,
    'delegation-stakeholder': 11,
    'ascend-plus': 12,
    'conflict-performance': 13,
    'people-leader-academy': 14,
    'coaching-next-line': 15,
    'leading-change-scale': 16,

    # Executive Courses (slides 18-27)
    'one-voice': 18,
    'enterprise-thinking': 19,
    'high-performance-culture': 20,
    'change-leadership-transformation': 21,
    'leadership-coaching-cross-border': 22,
    'enterprise-mindset-strategy': 23,
    'global-mobility': 24,
    'summit-innovation': 25,
    'miscellaneous': 26,
    'global-perspectives': 27,
}

CATEGORY_HEADER_SLIDES = {
    'individual-contributor': 2,
    'manager': 8,
    'executive': 17,
}

# Case study slide mapping (slides 28-36)
CASE_STUDY_TO_SLIDE_MAP = {
    'critical-thinking-gcc': 29,
    'time-to-productivity': 30,
    'granulearn-digital': 31,
    'functional-onboarding': 32,
    'change-management': 33,
    'culture-behaviors': 34,
    'lead-with-intent': 35,
    'storytelling': 36,
}

CASE_STUDY_COVER_SLIDE = 28


def customize_template_cover_slide(prs, client_name):
    """Customize the first slide from the template with client name"""
    # The template's first slide is already in the presentation
    # We'll modify it to add the client name
    if len(prs.slides) > 0:
        cover_slide = prs.slides[0]

        # Add client name text box if client name is provided
        if client_name:
            # Add a text box for client name (positioning may need adjustment based on template)
            client_box = cover_slide.shapes.add_textbox(
                Inches(0.5), Inches(5.2), Inches(9), Inches(0.6)
            )
            client_frame = client_box.text_frame
            client_frame.text = f"Prepared for: {client_name}"
            client_para = client_frame.paragraphs[0]
            client_para.font.size = Pt(24)
            client_para.font.bold = True
            client_para.font.color.rgb = RGBColor(248, 181, 12)  # Gold
            client_para.alignment = PP_ALIGN.CENTER

    return cover_slide if len(prs.slides) > 0 else None


def create_config_summary_slide(prs, config):
    """Create configuration summary slide
    ALWAYS adds slide at the END of the presentation.
    """
    # Use a title-only layout or blank layout
    slide_layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    # Clear any existing placeholders
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            sp = shape.element
            sp.getparent().remove(sp)

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Configuration Summary"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)  # White text
    title_para.alignment = PP_ALIGN.LEFT

    # Configuration items
    y_pos = 1.8

    # Format pricing text
    pricing = config.get('pricing')
    pricing_text = 'Not specified'
    if pricing:
        if pricing.get('type') == 'fixed':
            amount = pricing.get('amount', 0)
            pricing_text = f"${amount:,.2f}"
        elif pricing.get('type') == 'per-head':
            price_per_head = pricing.get('pricePerHead', 0)
            min_employees = pricing.get('minimumEmployees', 0)
            pricing_text = f"${price_per_head:,.2f} per employee (minimum {min_employees} employees)"

    config_items = [
        ('Client Name', config.get('clientName', 'Not specified')),
        ('Stage', config.get('stage', {}).get('name', 'Not selected')),
        ('Ambition', config.get('ambition', {}).get('name', 'Not selected')),
        ('Facilitation Model', config.get('facilitation', '').capitalize()),
        ('Delivery Modality', config.get('modality', '').capitalize()),
        ('Training Recipes', ', '.join([r['name'] for r in config.get('recipes', [])])),
        ('Pricing', pricing_text),
    ]

    for label, value in config_items:
        # Label
        label_box = slide.shapes.add_textbox(
            Inches(1), Inches(y_pos), Inches(2.5), Inches(0.5)
        )
        label_frame = label_box.text_frame
        label_frame.text = f"{label}:"
        label_frame.word_wrap = False
        label_para = label_frame.paragraphs[0]
        label_para.font.size = Pt(22)
        label_para.font.bold = True
        label_para.font.color.rgb = RGBColor(248, 181, 12)  # Gold (matches template accent)

        # Value
        value_box = slide.shapes.add_textbox(
            Inches(3.7), Inches(y_pos), Inches(5.5), Inches(0.5)
        )
        value_frame = value_box.text_frame
        value_frame.text = value
        value_frame.word_wrap = True
        value_para = value_frame.paragraphs[0]
        value_para.font.size = Pt(18)
        value_para.font.color.rgb = RGBColor(255, 255, 255)  # White text

        y_pos += 0.7

    return slide


def create_executive_summary_slide(prs, executive_summary):
    """Create executive summary slide(s) - splits content across multiple slides if needed
    ALWAYS adds slides at the END of the presentation.
    """
    # Split executive summary into chunks that fit on one slide
    # Estimate: ~1800 characters per slide with 14pt font
    MAX_CHARS_PER_SLIDE = 1800

    # Split by paragraphs first
    paragraphs = executive_summary.split('\n\n')

    slides_content = []
    current_slide_text = []
    current_char_count = 0

    for para_text in paragraphs:
        if not para_text.strip():
            continue

        para_length = len(para_text)

        # If adding this paragraph would exceed the limit, start a new slide
        if current_char_count + para_length > MAX_CHARS_PER_SLIDE and current_slide_text:
            slides_content.append(current_slide_text)
            current_slide_text = [para_text.strip()]
            current_char_count = para_length
        else:
            current_slide_text.append(para_text.strip())
            current_char_count += para_length

    # Add the last slide's content
    if current_slide_text:
        slides_content.append(current_slide_text)

    # Create slides - ALWAYS at the end
    created_slides = []
    for slide_num, paragraphs_for_slide in enumerate(slides_content):
        # Use a title-only layout or blank layout
        slide_layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)

        # Clear any existing placeholders
        for shape in list(slide.shapes):
            if shape.is_placeholder:
                sp = shape.element
                sp.getparent().remove(sp)

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        if slide_num == 0:
            title_frame.text = "Executive Summary"
        else:
            title_frame.text = f"Executive Summary (continued)"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)  # White text
        title_para.alignment = PP_ALIGN.LEFT

        # Summary content
        content_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.5), Inches(8.4), Inches(4.5)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True

        for i, para_text in enumerate(paragraphs_for_slide):
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()

            p.text = para_text
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(255, 255, 255)  # White text
            p.level = 0
            p.space_after = Pt(10)

        created_slides.append(slide)

    return created_slides


def copy_slide_shapes(source_slide, target_slide):
    """Copy all shapes from source slide to target slide"""
    from copy import deepcopy

    for shape in source_slide.shapes:
        el = shape.element
        newel = deepcopy(el)
        target_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    return target_slide


@app.route('/generate-pptx', methods=['POST'])
def generate_pptx():
    """Generate customized PowerPoint presentation"""
    try:
        data = request.json
        config = data.get('config', {})
        executive_summary = data.get('executiveSummary', '')

        # Debug: Print the raw recipes data
        print(f"DEBUG: Raw recipes from config: {config.get('recipes', [])}")
        print(f"DEBUG: Raw case tiles from config: {config.get('caseTiles', [])}")

        selected_recipe_ids = [r['id'] for r in config.get('recipes', [])]
        selected_case_study_ids = [c['id'] for c in config.get('caseTiles', [])]
        print(f"DEBUG: Selected recipe IDs: {selected_recipe_ids}")
        print(f"DEBUG: Selected case study IDs: {selected_case_study_ids}")

        # Load the template (only to read from)
        template_path = os.path.join(
            os.path.dirname(__file__),
            '..',
            'ppt',
            'Jade 2025-12-18_ANSR template - 1 pagers.pptx'
        )

        if not os.path.exists(template_path):
            return jsonify({'error': 'Template file not found'}), 404

        # APPROACH: Create NEW presentation and build slides in correct order
        # NO deletion, NO moving - just add slides in the order we want them
        # Order: Cover, Exec Summary, Config Summary, Recipes, Case Studies

        template_prs = Presentation(template_path)

        # Create NEW blank presentation with same slide dimensions
        prs = Presentation()
        prs.slide_width = template_prs.slide_width
        prs.slide_height = template_prs.slide_height

        # Get blank layout
        blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

        print("DEBUG: Building presentation in correct order...")

        # 1. Add cover slide (copy from template slide 0)
        print("DEBUG: Adding cover slide")
        cover_slide = prs.slides.add_slide(blank_layout)
        copy_slide_shapes(template_prs.slides[0], cover_slide)

        # Customize cover with client name
        if config.get('clientName'):
            client_box = cover_slide.shapes.add_textbox(
                Inches(0.5), Inches(5.2), Inches(9), Inches(0.6)
            )
            client_frame = client_box.text_frame
            client_frame.text = f"Prepared for: {config.get('clientName')}"
            client_para = client_frame.paragraphs[0]
            client_para.font.size = Pt(24)
            client_para.font.bold = True
            client_para.font.color.rgb = RGBColor(248, 181, 12)
            client_para.alignment = PP_ALIGN.CENTER

        # 2. Add executive summary slides (custom created)
        if executive_summary:
            print("DEBUG: Adding executive summary slides")
            create_executive_summary_slide(prs, executive_summary)

        # 3. Add config summary slide (custom created)
        print("DEBUG: Adding config summary slide")
        create_config_summary_slide(prs, config)

        # 4. Add selected recipe slides (copy from template)
        if selected_recipe_ids:
            print(f"DEBUG: Adding {len(selected_recipe_ids)} recipe slides")
            for recipe_id in selected_recipe_ids:
                slide_num = RECIPE_TO_SLIDE_MAP.get(recipe_id)
                if slide_num:
                    template_slide_idx = slide_num - 1
                    recipe_slide = prs.slides.add_slide(blank_layout)
                    copy_slide_shapes(template_prs.slides[template_slide_idx], recipe_slide)
                    print(f"DEBUG: Copied recipe slide {slide_num}")

        # 5. Add case studies if selected
        if selected_case_study_ids:
            # Add case studies cover slide
            print("DEBUG: Adding case studies cover slide")
            case_cover_slide = prs.slides.add_slide(blank_layout)
            copy_slide_shapes(template_prs.slides[CASE_STUDY_COVER_SLIDE - 1], case_cover_slide)

            # Add individual case study slides
            print(f"DEBUG: Adding {len(selected_case_study_ids)} case study slides")
            for case_study_id in selected_case_study_ids:
                slide_num = CASE_STUDY_TO_SLIDE_MAP.get(case_study_id)
                if slide_num:
                    template_slide_idx = slide_num - 1
                    case_slide = prs.slides.add_slide(blank_layout)
                    copy_slide_shapes(template_prs.slides[template_slide_idx], case_slide)
                    print(f"DEBUG: Copied case study slide {slide_num}")

        print(f"DEBUG: Final: {len(prs.slides)} slides")
        print(f"DEBUG: Order: Cover, Recipes, Case Studies, Exec Summary, Config Summary")
        print(f"DEBUG: NOTE: Exec/Config at end - manual reorder in PowerPoint needed")

        # Save to BytesIO
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)

        # Generate filename
        stage_id = config.get('stage', {}).get('id', 'config')
        date_str = datetime.now().strftime('%Y-%m-%d')
        filename = f"LD_Pitch_{stage_id}_{date_str}.pptx"

        return send_file(
            output,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name=filename
        )

    except Exception as e:
        print(f"Error generating PPTX: {str(e)}")
        import traceback
        traceback.print_exc()
        return jsonify({'error': str(e)}), 500


@app.route('/generate-summary', methods=['POST'])
def generate_summary():
    """
    Generate executive summary using Gemini API (secure backend)

    Request JSON should contain the configuration state from frontend
    Returns: {"summary": "Generated executive summary text..."}
    """
    try:
        if not GEMINI_API_KEY:
            return jsonify({'error': 'Gemini API key not configured on server'}), 500

        # Get configuration from request
        config = request.json

        if not config:
            return jsonify({'error': 'No configuration provided'}), 400

        # Build the prompt (matching frontend logic)
        client_name = config.get('clientName', 'the client')
        stage = config.get('stage') or {}
        ambition = config.get('ambition') or {}
        path = config.get('path') or {}
        facilitation = config.get('facilitation', '')
        modality = config.get('modality', '')
        recipes = config.get('recipes', [])
        case_tiles = config.get('caseTiles', [])

        # Format recipes and case tiles for prompt
        recipes_text = '\n'.join([f"- {r.get('name', '')}: {r.get('description', '')}" for r in recipes])
        case_tiles_text = '\n'.join([f"- {c.get('title', '')}: {c.get('metric', '')}" for c in case_tiles])

        path_type = 'Certification-Based' if path.get('type') == 'certification' else 'Tailored Programs'

        prompt = f"""
You are an expert L&D consultant creating an executive summary for a customized Learning & Development proposal for a Global Capability Center (GCC).

Based on the following configuration, write a compelling executive summary (3-4 paragraphs, approximately 250-300 words):

**Client Configuration:**
- Client Name: {client_name}
- Stage: {stage.get('name', '')} - {stage.get('description', '')}
- Strategic Ambition: {ambition.get('name', '')}
- Learning Path Type: {path_type}
- Facilitation Model: {facilitation}
- Delivery Modality: {modality}

**Selected Training Recipes ({len(recipes)}):**
{recipes_text}

**Selected Case Studies ({len(case_tiles)}):**
{case_tiles_text}

Write an executive summary that:
1. Opens with the strategic context and the client's current stage
2. Articulates the strategic ambition and how this L&D approach supports it
3. Highlights the selected training recipes and their expected impact
4. References the proven results from case studies
5. Concludes with the delivery approach and next steps

{f'IMPORTANT: When referring to the client, use "{client_name}" instead of placeholders like "[GCC Name]" or generic terms.' if client_name and client_name != 'the client' else ''}

Use a professional, consultative tone. Focus on business outcomes and strategic value. Make it compelling and actionable.

IMPORTANT: Output ONLY the executive summary paragraphs. Do NOT include any preamble text like "Here's an executive summary..." or headings like "Executive Summary:". Start directly with the content.
"""

        # Call Gemini API
        print("Generating executive summary with Gemini API...")
        model = genai.GenerativeModel('gemini-2.0-flash-exp')
        response = model.generate_content(prompt)

        if not response.text:
            return jsonify({'error': 'No summary generated'}), 500

        print("Successfully generated executive summary")
        return jsonify({'summary': response.text})

    except Exception as e:
        print(f"Error generating executive summary: {str(e)}")
        import traceback
        traceback.print_exc()
        return jsonify({'error': f'Failed to generate summary: {str(e)}'}), 500


@app.route('/health', methods=['GET'])
def health():
    """Health check endpoint"""
    return jsonify({'status': 'ok'})


if __name__ == '__main__':
    app.run(debug=True, port=5001)
